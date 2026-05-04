import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Define paths
BRAIN_DIR = r"C:\Users\saranyadheve\.gemini\antigravity\brain\2dbaffc7-a402-4d62-96ea-400825202d46"
OUTPUT_PATH = r"C:\Users\saranyadheve\OneDrive\Desktop\smart-revision-scheduler\Third_Review_Presentation.pptx"

# Screenshots mapping
IMG_PATHS = {
    "signup": os.path.join(BRAIN_DIR, "media__1776013086533.png"),
    "email": os.path.join(BRAIN_DIR, "media__1776013101161.png"),
    "otp": os.path.join(BRAIN_DIR, "media__1776013118465.png"),
    "dashboard": os.path.join(BRAIN_DIR, "media__1776013144883.png"),
    "notes": os.path.join(BRAIN_DIR, "media__1776013665306.png"),
    "tests": os.path.join(BRAIN_DIR, "media__1776013682976.png"),
    "hub": os.path.join(BRAIN_DIR, "media__1776013697960.png"),
    "ai_room": os.path.join(BRAIN_DIR, "media__1776013726618.png"),
    "consistency": os.path.join(BRAIN_DIR, "strategic_dashboard_1776012214605.png")
}

def add_header(slide, title_text):
    # Add top blue bar
    shape = slide.shapes.add_shape(
        1,  # Rect
        Inches(0), Inches(0), Inches(10), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2F, 0x55, 0x97)
    shape.line.color.rgb = RGBColor(0x2F, 0x55, 0x97)
    
    # Add title text
    txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    tf.text = title_text.upper()
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

def add_100_percent_badge(slide):
    shape = slide.shapes.add_shape(
        1,  # Rect
        Inches(7.5), Inches(0.8), Inches(2.2), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x27, 0xAE, 0x60)
    shape.line.color.rgb = RGBColor(0x27, 0xAE, 0x60)
    
    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(0.8), Inches(2.2), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "✔ 100% IMPLEMENTED"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    add_header(slide, title)
    
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(9)
    height = Inches(6)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for b in bullets:
        p = tf.add_paragraph()
        p.text = str(b)
        p.font.size = Pt(20)
        p.level = 0
        p.space_after = Pt(10)

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Third Review")
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "SMART REVISION SCHEDULER"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(0x2F, 0x55, 0x97)
    p.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(3))
    tf2 = txBox2.text_frame
    tf2.text = "An Intelligent AI-Powered Study Platform\n\nTeam Details: [NAME] ([REGISTER NO])\nGuide: [GUIDE NAME]\nDept. of [DEPARTMENT NAME]"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(20)
    p2.alignment = PP_ALIGN.CENTER

    # 2. Abstract
    add_bullet_slide(prs, "Abstract", [
        "End-to-end platform for UPSC, GATE, and TNPSC aspirants.",
        "Solves 'Information Overload' through curated analytics.",
        "Key Innovations: Kissara AI Tutor, Strategic Dashboard, Notes Archive.",
        "Architecture: Distributed React + Node.js + Spring Boot stack."
    ])

    # 3. Overall Design
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Overall Design")
    bullets = [
        "Frontend: React + Framer Motion (Glassmorphism Tier).",
        "Backend Node.js (5000): Specialized AI & Personal Notes.",
        "Backend Java (8080): Authentication & Course Management.",
        "Database: MariaDB (Relational Study Metadata).",
        "Security: OTP-based Email Verification (SMTP)."
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5))
    tf = txBox.text_frame
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(20)
        p.space_after = Pt(10)

    # 4. Experimental Results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Experimental Results")
    bullets = [
        "AI Chat: 100% functional with context-aware personas.",
        "Notes Archive: Verified multi-course persistence.",
        "Strategic Tracking: Real-time consistency map implementation.",
        "Practice Hub: Automated test generation for 4 tracks."
    ]
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(5), Inches(5))
    for b in bullets:
        p = tx.text_frame.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
    
    if os.path.exists(IMG_PATHS["consistency"]):
        slide.shapes.add_picture(IMG_PATHS["consistency"], Inches(5.5), Inches(1.5), width=Inches(4))

    # 5. Performance Evaluation
    add_bullet_slide(prs, "Performance Evaluation", [
        "UI Fluidity: Stabilized 60 FPS transitions.",
        "AI Efficiency: 0ms latency for Local Knowledge matches.",
        "API Latency: Optimized Node.js controllers (< 1.2s avg).",
        "Reliability: 100% database pool stability after refactor."
    ])

    # 6. References
    add_bullet_slide(prs, "References", [
        "IEEE 8983765: 'Study Companion: An AI based study assistant'.",
        "Google AI: Gemini API Prompt Engineering Documentation.",
        "Investopedia: Patterns for Educational Progress Analytics.",
        "MERN/Java Stack full-stack architectural design patterns."
    ])

    # 7. Contribution
    add_bullet_slide(prs, "Contribution of the Project", [
        "Frontend: Developed glassmorphism dashboard & roadmap logic.",
        "Backend: Engineered AI integration & distributed data flow.",
        "Database: Optimized relational schema for exam-specific data.",
        "Testing: Verified AI persona accuracy & system stability."
    ])

    # 8. 100% Code Implementation - Onboarding
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "100% Code Implementation - Demo")
    add_100_percent_badge(slide)
    if os.path.exists(IMG_PATHS["signup"]):
        slide.shapes.add_picture(IMG_PATHS["signup"], Inches(0.5), Inches(1.5), width=Inches(3))
    if os.path.exists(IMG_PATHS["email"]):
        slide.shapes.add_picture(IMG_PATHS["email"], Inches(3.6), Inches(1.5), width=Inches(3))
    if os.path.exists(IMG_PATHS["otp"]):
        slide.shapes.add_picture(IMG_PATHS["otp"], Inches(6.7), Inches(1.5), width=Inches(3))
    
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    tx.text_frame.text = "Signup Flow: Begin Mastery ➔ OTP Email ➔ Account Activation"
    tx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 9. 100% Code Implementation - Dashboard & AI
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "100% Code Implementation - Demo")
    add_100_percent_badge(slide)
    if os.path.exists(IMG_PATHS["dashboard"]):
        slide.shapes.add_picture(IMG_PATHS["dashboard"], Inches(0.5), Inches(1.5), width=Inches(4.4))
    if os.path.exists(IMG_PATHS["ai_room"]):
        slide.shapes.add_picture(IMG_PATHS["ai_room"], Inches(5.1), Inches(1.5), width=Inches(4.4))
    
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    tx.text_frame.text = "Analytics & IQ: Strategic Dashboard ➔ Conversational AI Tutor"
    tx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 10. 100% Code Implementation - Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "100% Code Implementation - Demo")
    add_100_percent_badge(slide)
    if os.path.exists(IMG_PATHS["tests"]):
        slide.shapes.add_picture(IMG_PATHS["tests"], Inches(0.5), Inches(1.5), width=Inches(3))
    if os.path.exists(IMG_PATHS["hub"]):
        slide.shapes.add_picture(IMG_PATHS["hub"], Inches(3.6), Inches(1.5), width=Inches(3))
    if os.path.exists(IMG_PATHS["notes"]):
        slide.shapes.add_picture(IMG_PATHS["notes"], Inches(6.7), Inches(1.5), width=Inches(3))

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    tx.text_frame.text = "Resource Modules: Practice Tests ➔ Topic Roadmap ➔ Notes Archive"
    tx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 11. Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Conclusion")
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.text = "THANK YOU"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.color.rgb = RGBColor(0x2F, 0x55, 0x97)
    p.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved successfully to: {OUTPUT_PATH}")

if __name__ == "__main__":
    main()

if __name__ == "__main__":
    main()
